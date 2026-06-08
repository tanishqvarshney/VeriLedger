"""
VeriLedger — Pitch Deck Generator
Generates a professional .pptx with a premium dark slate theme.
Run: python generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Palette ──────────────────────────────────────────────────────────────────
BG_BASE        = RGBColor(0x0B, 0x0D, 0x10)   # #0B0D10
BG_SURFACE     = RGBColor(0x12, 0x15, 0x1C)   # #12151C
BG_CARD        = RGBColor(0x1A, 0x1E, 0x28)   # #1A1E28
ACCENT_EMERALD = RGBColor(0x30, 0xD1, 0x58)   # #30D158
ACCENT_BLUE    = RGBColor(0x0A, 0x84, 0xFF)   # #0A84FF
ACCENT_AMBER   = RGBColor(0xFF, 0x9F, 0x0A)   # #FF9F0A
ACCENT_ROSE    = RGBColor(0xFF, 0x45, 0x3A)   # #FF453A
TEXT_PRIMARY   = RGBColor(0xF5, 0xF5, 0xF7)   # #F5F5F7
TEXT_SECONDARY = RGBColor(0x8E, 0x8E, 0x93)   # #8E8E93
TEXT_TERTIARY  = RGBColor(0x4A, 0x4A, 0x4F)   # #4A4A4F
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)

# Slide dimensions: 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]  # Completely blank


def fill_bg(slide, color=BG_BASE):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, radius=False):
    """Add a filled rectangle (no border)."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=TEXT_PRIMARY,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = "Inter"
    return txb


def add_multiline(slide, lines, left, top, width, height,
                  font_size=14, color=TEXT_SECONDARY, spacing=1.15):
    """lines: list of (text, bold, color_override)"""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for (txt, bold, col) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = col or color
        run.font.name = "Inter"
    return txb


def accent_bar(slide, color=ACCENT_BLUE, top=Inches(0.0), height=Inches(0.06)):
    """Full-width horizontal accent bar."""
    add_rect(slide, 0, top, SLIDE_W, height, color)


def slide_label(slide, text, color=ACCENT_BLUE):
    """Small monospaced label in top-left."""
    add_text(slide, text,
             left=Inches(0.55), top=Inches(0.18),
             width=Inches(4), height=Inches(0.35),
             font_size=9, color=color, bold=True)


def footer(slide, text="VeriLedger  •  Confidential"):
    add_text(slide, text,
             left=Inches(0), top=Inches(7.1),
             width=SLIDE_W, height=Inches(0.3),
             font_size=8, color=TEXT_TERTIARY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    """Cover — title + tagline + accent geometry."""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide)

    # Bottom-right geometric accent block
    add_rect(slide, Inches(9.5), Inches(4.5), Inches(3.833), Inches(3.0), BG_SURFACE)
    add_rect(slide, Inches(11.0), Inches(5.8), Inches(2.333), Inches(1.7), BG_CARD)

    # Top accent line
    accent_bar(slide, ACCENT_EMERALD, top=Inches(0.0), height=Inches(0.05))

    # Logo / wordmark
    add_text(slide, "VERILEDGER",
             left=Inches(0.7), top=Inches(2.0),
             width=Inches(9), height=Inches(1.2),
             font_size=72, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Emerald underline block
    add_rect(slide, Inches(0.7), Inches(3.18), Inches(1.6), Inches(0.08), ACCENT_EMERALD)

    # Tagline
    add_text(slide, "Real-Time Document Forgery Detection\n& Cryptographic Audit Layer for Bank Underwriting",
             left=Inches(0.7), top=Inches(3.38),
             width=Inches(8.5), height=Inches(1.1),
             font_size=20, bold=False, color=TEXT_SECONDARY, align=PP_ALIGN.LEFT)

    # Theme badge
    badge = add_rect(slide, Inches(0.7), Inches(4.7), Inches(4.2), Inches(0.42), BG_CARD)
    add_text(slide, "🏆  Agentic Regulatory Intelligence & Compliance",
             left=Inches(0.75), top=Inches(4.72),
             width=Inches(4.1), height=Inches(0.38),
             font_size=11, bold=True, color=ACCENT_EMERALD, align=PP_ALIGN.LEFT)

    # Stats row
    for i, (val, lbl) in enumerate([("100%", "Test Pass Rate"), ("11", "Forensic Tests"), ("3", "Forensic Engines")]):
        x = Inches(0.7 + i * 2.1)
        add_text(slide, val, left=x, top=Inches(5.5), width=Inches(2), height=Inches(0.55),
                 font_size=28, bold=True, color=ACCENT_EMERALD)
        add_text(slide, lbl, left=x, top=Inches(6.05), width=Inches(2), height=Inches(0.35),
                 font_size=10, color=TEXT_SECONDARY)

    footer(slide)


def slide_02_problem(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide)
    accent_bar(slide, ACCENT_ROSE)
    slide_label(slide, "01  /  PROBLEM", ACCENT_ROSE)

    add_text(slide, "Loan Fraud Costs Banks Billions",
             left=Inches(0.55), top=Inches(0.6),
             width=Inches(10), height=Inches(0.7),
             font_size=36, bold=True, color=WHITE)

    add_text(slide, "Manual underwriting review cannot keep pace with increasingly sophisticated document forgery.",
             left=Inches(0.55), top=Inches(1.35),
             width=Inches(9.5), height=Inches(0.5),
             font_size=15, color=TEXT_SECONDARY)

    # Stat cards
    stats = [
        ("$3.1B+",  "lost to mortgage & personal loan fraud annually (US, 2023)", ACCENT_ROSE),
        ("~6 min",  "average manual review time per document — impossible at scale", ACCENT_AMBER),
        ("78%",     "of fraudulent documents pass unaided visual inspection", ACCENT_AMBER),
        ("0",       "cryptographic proof that a reviewed doc was actually scanned", TEXT_SECONDARY),
    ]
    for i, (stat, desc, col) in enumerate(stats):
        col_idx = i % 2
        row_idx = i // 2
        x = Inches(0.55 + col_idx * 6.3)
        y = Inches(2.1  + row_idx * 2.1)
        add_rect(slide, x, y, Inches(5.9), Inches(1.85), BG_SURFACE)
        add_rect(slide, x, y, Inches(0.06), Inches(1.85), col)
        add_text(slide, stat, left=x+Inches(0.2), top=y+Inches(0.15),
                 width=Inches(5.5), height=Inches(0.6),
                 font_size=32, bold=True, color=col)
        add_text(slide, desc, left=x+Inches(0.2), top=y+Inches(0.7),
                 width=Inches(5.5), height=Inches(1.0),
                 font_size=12, color=TEXT_SECONDARY)

    footer(slide)


def slide_03_solution(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide)
    accent_bar(slide, ACCENT_EMERALD)
    slide_label(slide, "02  /  SOLUTION", ACCENT_EMERALD)

    add_text(slide, "VeriLedger",
             left=Inches(0.55), top=Inches(0.6),
             width=Inches(7), height=Inches(0.7),
             font_size=40, bold=True, color=WHITE)

    add_text(slide, "An automated forensic document trust console that ingests any uploaded PDF or image,\n"
                    "runs three independent forensic engines, and returns a Trust Score with full visual & metadata evidence.",
             left=Inches(0.55), top=Inches(1.38),
             width=Inches(12.2), height=Inches(0.75),
             font_size=13.5, color=TEXT_SECONDARY)

    # 3 Engine cards
    engines = [
        ("🔬", "Error Level\nAnalysis", "Detects pixel-grid compression anomalies — precisely locates spliced or altered numbers.", ACCENT_BLUE),
        ("🧬", "Copy-Move\nDetection", "ORB keypoint matching exposes cloned stamps, forged signatures, or pasted text blocks.", ACCENT_EMERALD),
        ("📋", "Metadata &\nStructure Audit", "Flags suspicious PDF creators (iLovePDF, Photoshop), date-time deltas, and font issues.", ACCENT_AMBER),
    ]
    for i, (icon, title, desc, col) in enumerate(engines):
        x = Inches(0.55 + i * 4.25)
        y = Inches(2.35)
        add_rect(slide, x, y, Inches(4.0), Inches(2.8), BG_SURFACE)
        add_rect(slide, x, y, Inches(4.0), Inches(0.07), col)
        add_text(slide, icon, left=x+Inches(0.2), top=y+Inches(0.22),
                 width=Inches(0.6), height=Inches(0.5), font_size=28, color=WHITE)
        add_text(slide, title, left=x+Inches(0.2), top=y+Inches(0.72),
                 width=Inches(3.6), height=Inches(0.65),
                 font_size=16, bold=True, color=WHITE)
        add_text(slide, desc, left=x+Inches(0.2), top=y+Inches(1.35),
                 width=Inches(3.6), height=Inches(1.3),
                 font_size=11.5, color=TEXT_SECONDARY)

    # Plus two more features
    extras = [
        ("🔗", "Cryptographic Ledger", "Every document hashed into a SHA-256 chained SQLite ledger — tamper-evident audit trail.", ACCENT_BLUE),
        ("📊", "Cross-Doc OCR Reconciliation", "Salary slip vs. bank statement — automatically flags name & deposit mismatches.", ACCENT_EMERALD),
    ]
    for i, (icon, title, desc, col) in enumerate(extras):
        x = Inches(0.55 + i * 6.4)
        y = Inches(5.4)
        add_rect(slide, x, y, Inches(6.1), Inches(1.65), BG_CARD)
        add_rect(slide, x, y, Inches(0.06), Inches(1.65), col)
        add_text(slide, f"{icon}  {title}", left=x+Inches(0.2), top=y+Inches(0.15),
                 width=Inches(5.7), height=Inches(0.45),
                 font_size=14, bold=True, color=WHITE)
        add_text(slide, desc, left=x+Inches(0.2), top=y+Inches(0.6),
                 width=Inches(5.7), height=Inches(0.9),
                 font_size=11.5, color=TEXT_SECONDARY)

    footer(slide)


def slide_04_how_it_works(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide)
    accent_bar(slide, ACCENT_BLUE)
    slide_label(slide, "03  /  HOW IT WORKS", ACCENT_BLUE)

    add_text(slide, "End-to-End Pipeline",
             left=Inches(0.55), top=Inches(0.6),
             width=Inches(10), height=Inches(0.65),
             font_size=34, bold=True, color=WHITE)

    steps = [
        ("01", "Document\nIngested", "PDF or image uploaded via drag-and-drop. SHA-256 hash computed immediately.", ACCENT_BLUE),
        ("02", "Forensic\nEngines Run", "ELA, Copy-Move, and Metadata audit run in parallel on the rasterized document.", ACCENT_EMERALD),
        ("03", "Trust Score\nFused", "Weighted formula: 35% ELA + 35% Clone + 30% Metadata = Trust Score 0–100.", ACCENT_AMBER),
        ("04", "Ledger Block\nCommitted", "Block (doc_hash → prev_hash → chain_hash) appended to tamper-evident DB.", ACCENT_ROSE),
        ("05", "Dashboard\nUpdated", "Radial gauge, ELA heatmap, clone map, and ledger timeline rendered live.", TEXT_SECONDARY),
    ]

    for i, (num, title, desc, col) in enumerate(steps):
        x = Inches(0.3 + i * 2.57)
        y = Inches(1.6)

        # Card
        add_rect(slide, x, y, Inches(2.4), Inches(4.5), BG_SURFACE)
        add_rect(slide, x, y, Inches(2.4), Inches(0.06), col)

        # Step number
        add_text(slide, num, left=x+Inches(0.15), top=y+Inches(0.2),
                 width=Inches(0.7), height=Inches(0.45),
                 font_size=22, bold=True, color=col)

        # Title
        add_text(slide, title, left=x+Inches(0.15), top=y+Inches(0.7),
                 width=Inches(2.1), height=Inches(0.8),
                 font_size=14, bold=True, color=WHITE)

        # Desc
        add_text(slide, desc, left=x+Inches(0.15), top=y+Inches(1.55),
                 width=Inches(2.1), height=Inches(2.6),
                 font_size=11, color=TEXT_SECONDARY)

        # Arrow between cards
        if i < len(steps) - 1:
            add_text(slide, "→",
                     left=x+Inches(2.4), top=y+Inches(1.9),
                     width=Inches(0.18), height=Inches(0.5),
                     font_size=18, color=TEXT_TERTIARY, align=PP_ALIGN.CENTER)

    footer(slide)


def slide_05_trust_score(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide)
    accent_bar(slide, ACCENT_EMERALD)
    slide_label(slide, "04  /  TRUST SCORE ENGINE", ACCENT_EMERALD)

    add_text(slide, "How Trust Is Calculated",
             left=Inches(0.55), top=Inches(0.6),
             width=Inches(10), height=Inches(0.65),
             font_size=34, bold=True, color=WHITE)

    # Formula card
    add_rect(slide, Inches(0.55), Inches(1.5), Inches(12.2), Inches(1.2), BG_SURFACE)
    add_text(slide, "Trust Score  =  100  –  (ELA Risk × 0.35)  –  (Clone Risk × 0.35)  –  (Metadata Penalty × 0.30)",
             left=Inches(0.75), top=Inches(1.65),
             width=Inches(12.0), height=Inches(0.8),
             font_size=16, bold=True, color=ACCENT_EMERALD, align=PP_ALIGN.CENTER)

    # Score examples
    examples = [
        ("95–100", "TRUSTED",   "Clean salary slip — no ELA spikes, zero clone matches, no suspicious metadata.", ACCENT_EMERALD),
        ("50–75",  "REVIEW",    "Mixed signals — minor ELA variance or a single unembedded font. Requires manual check.", ACCENT_AMBER),
        ("0–35",   "REJECTED",  "Tampered document — Photoshop creator string, cloned stamp blocks, bright ELA hotspot at Net Pay.", ACCENT_ROSE),
    ]
    for i, (score, verdict, desc, col) in enumerate(examples):
        x = Inches(0.55 + i * 4.25)
        y = Inches(2.95)
        add_rect(slide, x, y, Inches(4.0), Inches(2.9), BG_CARD)
        add_rect(slide, x, y, Inches(4.0), Inches(0.07), col)
        add_text(slide, score, left=x+Inches(0.2), top=y+Inches(0.2),
                 width=Inches(3.6), height=Inches(0.65),
                 font_size=36, bold=True, color=col)
        add_text(slide, verdict, left=x+Inches(0.2), top=y+Inches(0.85),
                 width=Inches(3.6), height=Inches(0.4),
                 font_size=13, bold=True, color=WHITE)
        add_text(slide, desc, left=x+Inches(0.2), top=y+Inches(1.3),
                 width=Inches(3.6), height=Inches(1.4),
                 font_size=11, color=TEXT_SECONDARY)

    # Weights row
    add_text(slide, "Signal Weights:",
             left=Inches(0.55), top=Inches(6.1),
             width=Inches(2.5), height=Inches(0.38),
             font_size=11, bold=True, color=TEXT_SECONDARY)
    for i, (label, pct, col) in enumerate([("ELA", "35%", ACCENT_BLUE), ("Copy-Move", "35%", ACCENT_EMERALD), ("Metadata", "30%", ACCENT_AMBER)]):
        x = Inches(2.9 + i * 3.4)
        add_rect(slide, x, Inches(6.15), Inches(3.0), Inches(0.3), BG_SURFACE)
        add_rect(slide, x, Inches(6.15), Inches(3.0 * float(pct[:-1]) / 100), Inches(0.3), col)
        add_text(slide, f"{label}  {pct}", left=x+Inches(0.1), top=Inches(6.5),
                 width=Inches(2.9), height=Inches(0.35),
                 font_size=10, color=col, bold=True)

    footer(slide)


def slide_06_ledger(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide)
    accent_bar(slide, ACCENT_BLUE)
    slide_label(slide, "05  /  CRYPTOGRAPHIC LEDGER", ACCENT_BLUE)

    add_text(slide, "Tamper-Evident Audit Chain",
             left=Inches(0.55), top=Inches(0.6),
             width=Inches(10), height=Inches(0.65),
             font_size=34, bold=True, color=WHITE)

    add_text(slide, "Every document processed by VeriLedger is permanently registered in a cryptographic SQLite block chain.",
             left=Inches(0.55), top=Inches(1.38),
             width=Inches(12.2), height=Inches(0.45),
             font_size=13, color=TEXT_SECONDARY)

    # Chain diagram
    blocks = [
        ("Block #1", "Genesis",           "prev: 0000…0000"),
        ("Block #2", "clean_salary.pdf",  "prev: a3f9…c82e"),
        ("Block #3", "tampered_doc.pdf",  "prev: 7b21…d94a"),
        ("Block #N", "…next document",    "prev: e19c…f03b"),
    ]
    for i, (num, name, prev) in enumerate(blocks):
        x = Inches(0.45 + i * 3.25)
        y = Inches(2.15)
        add_rect(slide, x, y, Inches(3.0), Inches(2.5), BG_SURFACE)
        add_rect(slide, x, y, Inches(3.0), Inches(0.07), ACCENT_BLUE)
        add_text(slide, num,  left=x+Inches(0.15), top=y+Inches(0.18),
                 width=Inches(2.7), height=Inches(0.35), font_size=10, bold=True, color=ACCENT_BLUE)
        add_text(slide, name, left=x+Inches(0.15), top=y+Inches(0.55),
                 width=Inches(2.7), height=Inches(0.45), font_size=13, bold=True, color=WHITE)
        add_text(slide, prev, left=x+Inches(0.15), top=y+Inches(1.05),
                 width=Inches(2.7), height=Inches(0.35), font_size=9, color=TEXT_TERTIARY)
        add_text(slide, "chain_hash: SHA-256\n(doc_hash + prev_hash)", left=x+Inches(0.15), top=y+Inches(1.5),
                 width=Inches(2.7), height=Inches(0.75), font_size=9, color=TEXT_SECONDARY)
        if i < len(blocks) - 1:
            add_text(slide, "→", left=x+Inches(3.0), top=y+Inches(1.0),
                     width=Inches(0.25), height=Inches(0.5), font_size=20, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    # Properties list
    props = [
        ("🔐  Immutable Chain",     "Any mutation to a past block breaks the chain_hash for ALL subsequent blocks."),
        ("🕵️  Audit Endpoint",     "GET /ledger/verify traverses every block and flags the exact mutated block ID."),
        ("📌  Zero Trust Storage", "Even if the DBA edits the DB directly, VeriLedger detects and reports the breach."),
    ]
    for i, (title, desc) in enumerate(props):
        y = Inches(5.0 + i * 0.65)
        add_rect(slide, Inches(0.55), y, Inches(12.2), Inches(0.55), BG_CARD)
        add_text(slide, title, left=Inches(0.75), top=y+Inches(0.08),
                 width=Inches(3.5), height=Inches(0.4), font_size=11, bold=True, color=WHITE)
        add_text(slide, desc, left=Inches(4.3), top=y+Inches(0.08),
                 width=Inches(8.3), height=Inches(0.4), font_size=11, color=TEXT_SECONDARY)

    footer(slide)


def slide_07_tech_stack(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide)
    accent_bar(slide, ACCENT_AMBER)
    slide_label(slide, "06  /  TECHNOLOGY STACK", ACCENT_AMBER)

    add_text(slide, "Built on Production-Grade Tooling",
             left=Inches(0.55), top=Inches(0.6),
             width=Inches(10), height=Inches(0.65),
             font_size=34, bold=True, color=WHITE)

    categories = [
        ("Backend", [
            ("FastAPI + Uvicorn",    "High-performance async REST API server"),
            ("Python 3.13",         "Core runtime with type-annotated handlers"),
            ("SQLite",              "Embedded cryptographic ledger database"),
        ], ACCENT_BLUE),
        ("Forensics", [
            ("OpenCV",              "ORB keypoints & brute-force KNN matching for copy-move detection"),
            ("Pillow (PIL)",        "JPEG re-compression for Error Level Analysis"),
            ("PyMuPDF (fitz)",      "PDF rasterization and metadata extraction"),
            ("Tesseract OCR",       "Text extraction for cross-document reconciliation"),
        ], ACCENT_EMERALD),
        ("Frontend", [
            ("React + Vite",        "Component-based SPA with HMR dev server"),
            ("Tailwind CSS v4",     "Utility-first dark-mode design system"),
            ("Lucide Icons",        "Consistent minimal icon set"),
        ], ACCENT_AMBER),
        ("DevOps", [
            ("Docker + Compose",    "Single-command containerized deployment"),
            ("pytest",              "11-test forensic discrimination suite"),
            ("GitHub",              "Version-controlled source: tanishqvarshney/VeriLedger"),
        ], ACCENT_ROSE),
    ]

    for i, (cat, items, col) in enumerate(categories):
        col_i = i % 2
        row_i = i // 2
        x = Inches(0.55 + col_i * 6.4)
        y = Inches(1.65 + row_i * 2.7)
        add_rect(slide, x, y, Inches(6.1), Inches(2.5), BG_SURFACE)
        add_rect(slide, x, y, Inches(6.1), Inches(0.07), col)
        add_text(slide, cat, left=x+Inches(0.2), top=y+Inches(0.15),
                 width=Inches(5.7), height=Inches(0.4), font_size=13, bold=True, color=col)
        for j, (tech, desc) in enumerate(items):
            ty = y + Inches(0.65 + j * 0.58)
            add_text(slide, f"• {tech}", left=x+Inches(0.2), top=ty,
                     width=Inches(2.3), height=Inches(0.48), font_size=11, bold=True, color=WHITE)
            add_text(slide, desc, left=x+Inches(2.55), top=ty,
                     width=Inches(3.35), height=Inches(0.48), font_size=10, color=TEXT_SECONDARY)

    footer(slide)


def slide_08_demo(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide)
    accent_bar(slide, ACCENT_EMERALD)
    slide_label(slide, "07  /  DEMO WALKTHROUGH", ACCENT_EMERALD)

    add_text(slide, "See It In Action",
             left=Inches(0.55), top=Inches(0.6),
             width=Inches(10), height=Inches(0.65),
             font_size=34, bold=True, color=WHITE)

    steps = [
        ("Step 1", "Upload Document", "Drag clean_salary_slip.pdf onto the upload zone. The scanner animation starts immediately.", ACCENT_BLUE, "Trust Score → 97.4  🟢"),
        ("Step 2", "Review ELA Heatmap", "Switch to ELA Split view. Drag the slider — the dark uniform heatmap confirms no splicing.", ACCENT_EMERALD, "ELA Risk → 2.1  ✅"),
        ("Step 3", "Upload Tampered Doc", "Drop tampered_salary_slip.pdf. Trust Score drops to 18.6 with a red glow. ELA shows bright hotspot on Net Pay.", ACCENT_ROSE, "Trust Score → 18.6  🔴"),
        ("Step 4", "Cross-Doc Audit", "Switch to Cross-Document tab. Upload both files — system flags Missing Salary Deposit: $5,400 not found.", ACCENT_AMBER, "Match Score → 55.0  ⚠️"),
        ("Step 5", "Ledger Integrity", "Click Audit Chain Integrity — green banner confirms all blocks intact. Mutate DB → red breach warning appears.", ACCENT_BLUE, "Chain: Intact  🔐"),
    ]

    for i, (step, title, desc, col, result) in enumerate(steps):
        y = Inches(1.5 + i * 1.1)
        add_rect(slide, Inches(0.55), y, Inches(12.2), Inches(0.95), BG_SURFACE if i % 2 == 0 else BG_CARD)
        add_rect(slide, Inches(0.55), y, Inches(0.06), Inches(0.95), col)
        add_text(slide, step, left=Inches(0.75), top=y+Inches(0.1),
                 width=Inches(0.8), height=Inches(0.35), font_size=9, bold=True, color=col)
        add_text(slide, title, left=Inches(1.55), top=y+Inches(0.1),
                 width=Inches(3.2), height=Inches(0.35), font_size=13, bold=True, color=WHITE)
        add_text(slide, desc, left=Inches(1.55), top=y+Inches(0.48),
                 width=Inches(8.2), height=Inches(0.35), font_size=10.5, color=TEXT_SECONDARY)
        add_text(slide, result, left=Inches(9.8), top=y+Inches(0.22),
                 width=Inches(2.8), height=Inches(0.45), font_size=12, bold=True, color=col, align=PP_ALIGN.RIGHT)

    footer(slide)


def slide_09_test_suite(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide)
    accent_bar(slide, ACCENT_EMERALD)
    slide_label(slide, "08  /  TEST SUITE", ACCENT_EMERALD)

    add_text(slide, "11 / 11 Tests Passing  ✓",
             left=Inches(0.55), top=Inches(0.6),
             width=Inches(12), height=Inches(0.65),
             font_size=34, bold=True, color=ACCENT_EMERALD)

    add_text(slide, "Every test verifies forensic discrimination — not just HTTP 200. Tests assert tampered ≠ clean.",
             left=Inches(0.55), top=Inches(1.38),
             width=Inches(12.2), height=Inches(0.45),
             font_size=13, color=TEXT_SECONDARY)

    tests = [
        ("ELA Discrimination",         "Tampered ELA score > Clean ELA score. Hotspot spatially near the edited amount — not uniform.", ACCENT_BLUE),
        ("Copy-Move Discrimination",   "Clone-stamp doc triggers detection. Clean text-heavy doc scores exactly 0.0 clone risk.", ACCENT_EMERALD),
        ("Metadata Flagging",          "iLovePDF / Photoshop creator strings flagged. ModDate > 180s delta flagged. Clean doc = 0 flags.", ACCENT_AMBER),
        ("Trust Score Weights",        "Mathematical weight validation: 0.35 ELA + 0.35 Clone + 0.30 Metadata. Deterministic & monotonic.", ACCENT_BLUE),
        ("Ledger Integrity Breach",    "3-doc chain built → row hash mutated in DB → /ledger/verify detects and reports mutated block ID.", ACCENT_ROSE),
        ("Cross-Doc OCR Reconcile",    "Name mismatch and Missing Salary Deposit ($5,400) flagged when comparing mismatched pair.", ACCENT_EMERALD),
        ("Error Handling",             "Unsupported types, empty files, oversized (>10MB), corrupted PDFs, encrypted files → HTTP 400.", ACCENT_AMBER),
    ]

    for i, (name, desc, col) in enumerate(tests):
        col_i = i % 2
        row_i = i // 2
        if i == 6:  # last card spans full width
            x = Inches(0.55)
            y = Inches(2.0 + 3 * 1.2)
            w = Inches(12.2)
        else:
            x = Inches(0.55 + col_i * 6.4)
            y = Inches(2.0 + row_i * 1.2)
            w = Inches(6.1)
        add_rect(slide, x, y, w, Inches(1.05), BG_SURFACE)
        add_rect(slide, x, y, Inches(0.06), Inches(1.05), col)
        add_text(slide, f"✓  {name}", left=x+Inches(0.15), top=y+Inches(0.1),
                 width=w-Inches(0.3), height=Inches(0.38), font_size=12, bold=True, color=WHITE)
        add_text(slide, desc, left=x+Inches(0.15), top=y+Inches(0.5),
                 width=w-Inches(0.3), height=Inches(0.45), font_size=10, color=TEXT_SECONDARY)

    footer(slide)


def slide_10_impact(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide)
    accent_bar(slide, ACCENT_BLUE)
    slide_label(slide, "09  /  IMPACT & ROADMAP", ACCENT_BLUE)

    add_text(slide, "From Prototype to Production",
             left=Inches(0.55), top=Inches(0.6),
             width=Inches(10), height=Inches(0.65),
             font_size=34, bold=True, color=WHITE)

    # Impact column
    add_text(slide, "Immediate Impact",
             left=Inches(0.55), top=Inches(1.5),
             width=Inches(6), height=Inches(0.45),
             font_size=16, bold=True, color=ACCENT_EMERALD)

    impacts = [
        "⚡️  Reduces per-document review time from ~6 min to < 10 seconds",
        "🛡️  Provides cryptographic proof that every document was forensically scanned",
        "📉  Detects tampering invisible to the naked eye via ELA heatmaps",
        "⚖️  Creates a compliance audit trail meeting AML/KYC audit requirements",
        "🔗  Immutable ledger blocks any retroactive alteration of the audit record",
    ]
    for i, imp in enumerate(impacts):
        add_rect(slide, Inches(0.55), Inches(2.05 + i * 0.82), Inches(5.9), Inches(0.68), BG_SURFACE)
        add_text(slide, imp, left=Inches(0.75), top=Inches(2.1 + i * 0.82),
                 width=Inches(5.5), height=Inches(0.55), font_size=12, color=TEXT_PRIMARY)

    # Roadmap column
    add_text(slide, "Roadmap",
             left=Inches(7.0), top=Inches(1.5),
             width=Inches(6), height=Inches(0.45),
             font_size=16, bold=True, color=ACCENT_BLUE)

    roadmap = [
        ("Q3 2026", "LLM-powered narrative fraud detector (semantic inconsistencies in salary descriptions)", ACCENT_BLUE),
        ("Q3 2026", "Real-time webhook integration with core banking APIs (Temenos, Finastra)", ACCENT_BLUE),
        ("Q4 2026", "Multi-jurisdiction compliance rule sets (RBI, FCA, OCC) embedded in audit engine", ACCENT_AMBER),
        ("Q4 2026", "Batch pipeline: scan 10,000 documents overnight, export CSV trust score report", ACCENT_AMBER),
        ("2027",    "Blockchain ledger migration (Hyperledger Fabric) for multi-bank consortium audits", ACCENT_EMERALD),
    ]
    for i, (quarter, item, col) in enumerate(roadmap):
        y = Inches(2.05 + i * 0.82)
        add_rect(slide, Inches(7.0), y, Inches(5.9), Inches(0.68), BG_CARD)
        add_rect(slide, Inches(7.0), y, Inches(0.06), Inches(0.68), col)
        add_text(slide, quarter, left=Inches(7.15), top=y+Inches(0.07),
                 width=Inches(0.9), height=Inches(0.3), font_size=9, bold=True, color=col)
        add_text(slide, item, left=Inches(8.1), top=y+Inches(0.07),
                 width=Inches(4.65), height=Inches(0.52), font_size=10.5, color=TEXT_PRIMARY)

    footer(slide)


def slide_11_closing(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide)

    # Full background accent geometry
    add_rect(slide, Inches(9.0), 0, Inches(4.333), SLIDE_H, BG_SURFACE)
    add_rect(slide, Inches(10.5), Inches(2.5), Inches(2.833), Inches(5.0), BG_CARD)
    accent_bar(slide, ACCENT_EMERALD, top=Inches(0.0), height=Inches(0.05))

    add_text(slide, "Trust Every Document.\nProtect Every Decision.",
             left=Inches(0.55), top=Inches(1.8),
             width=Inches(8.0), height=Inches(2.0),
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    add_rect(slide, Inches(0.55), Inches(3.9), Inches(1.8), Inches(0.07), ACCENT_EMERALD)

    add_text(slide, "VeriLedger brings institutional-grade forensic document analysis\nto every underwriting workflow — fast, precise, and audit-ready.",
             left=Inches(0.55), top=Inches(4.1),
             width=Inches(8.0), height=Inches(0.9),
             font_size=15, color=TEXT_SECONDARY)

    links = [
        ("GitHub", "github.com/tanishqvarshney/VeriLedger", ACCENT_BLUE),
        ("Live Demo", "localhost:3000", ACCENT_EMERALD),
    ]
    for i, (label, url, col) in enumerate(links):
        x = Inches(0.55 + i * 4.0)
        add_rect(slide, x, Inches(5.2), Inches(3.7), Inches(0.65), BG_CARD)
        add_text(slide, f"{label}:  {url}", left=x+Inches(0.2), top=Inches(5.28),
                 width=Inches(3.3), height=Inches(0.45), font_size=12, color=col, bold=True)

    add_text(slide, "VERILEDGER",
             left=Inches(9.2), top=Inches(3.0),
             width=Inches(4.0), height=Inches(1.0),
             font_size=32, bold=True, color=TEXT_TERTIARY, align=PP_ALIGN.CENTER)

    footer(slide, "VeriLedger  •  Agentic Regulatory Intelligence & Compliance  •  2026")


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()
    slide_01_cover(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_how_it_works(prs)
    slide_05_trust_score(prs)
    slide_06_ledger(prs)
    slide_07_tech_stack(prs)
    slide_08_demo(prs)
    slide_09_test_suite(prs)
    slide_10_impact(prs)
    slide_11_closing(prs)

    out = "VeriLedger_Pitch_Deck.pptx"
    prs.save(out)
    print(f"✅  Saved: {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    main()
